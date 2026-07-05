"""Document processing with text extraction and intelligent chunking"""
import io
import re
from typing import List, Tuple, Dict, Optional
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
import logging

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Extract and process text from various document formats"""
    
    SUPPORTED_FORMATS = {"pdf", "docx", "pptx", "txt", "csv", "md"}
    MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
    
    @staticmethod
    def validate_file(file_content: bytes, file_type: str) -> bool:
        """Validate file before processing"""
        if file_type.lower() not in DocumentProcessor.SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        if len(file_content) > DocumentProcessor.MAX_FILE_SIZE:
            raise ValueError(f"File too large: {len(file_content) / (1024*1024):.1f}MB")
        
        return True
    
    @staticmethod
    def extract_text(file_content: bytes, file_type: str) -> Tuple[str, dict]:
        """Extract text and metadata from document
        
        Args:
            file_content: Raw file bytes
            file_type: File extension (pdf, docx, etc)
        
        Returns:
            Tuple of (full_text, metadata)
        """
        file_type = file_type.lower().strip(".")
        
        try:
            DocumentProcessor.validate_file(file_content, file_type)
            
            if file_type == "pdf":
                return DocumentProcessor._extract_pdf(file_content)
            elif file_type == "docx":
                return DocumentProcessor._extract_docx(file_content)
            elif file_type == "pptx":
                return DocumentProcessor._extract_pptx(file_content)
            elif file_type == "txt":
                return DocumentProcessor._extract_txt(file_content)
            elif file_type == "csv":
                return DocumentProcessor._extract_csv(file_content)
            elif file_type == "md":
                return DocumentProcessor._extract_txt(file_content)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            raise
    
    @staticmethod
    def _extract_pdf(file_content: bytes) -> Tuple[str, dict]:
        """Extract text from PDF with page tracking and image handling
        
        Returns:
            Tuple of (full_text, {pages: int, text_by_page: dict})
        """
        try:
            reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text_by_page = {}
            empty_pages = []
            
            for page_num, page in enumerate(reader.pages, 1):
                try:
                    text = page.extract_text()
                    
                    # If page is mostly empty, it might have images
                    if not text or len(text.strip()) < 10:
                        logger.warning(f"Page {page_num}: Minimal text extracted (may be image-heavy)")
                        empty_pages.append(page_num)
                        # Set placeholder text
                        text = f"[Page {page_num} - Image or scanned content]"
                    
                    text_by_page[page_num] = text
                    
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num}: {e}")
                    empty_pages.append(page_num)
                    text_by_page[page_num] = f"[Page {page_num} - Extraction error: {str(e)}]"
            
            # Combine pages with page markers
            full_text_parts = []
            for page_num in sorted(text_by_page.keys()):
                full_text_parts.append(f"--- PAGE {page_num} ---\n{text_by_page[page_num]}\n")
            
            full_text = "\n".join(full_text_parts)
            
            # If all pages are empty/images, return minimal text
            if len(empty_pages) == len(reader.pages):
                logger.warning(f"PDF {len(reader.pages)} pages appear to be all images")
                full_text = f"PDF Document: {len(reader.pages)} pages (image-based content)"
            
            return full_text, {
                "pages": len(reader.pages),
                "text_by_page": text_by_page,
                "empty_pages": empty_pages,
                "format": "pdf"
            }
        except Exception as e:
            logger.error(f"PDF extraction error: {e}")
            # Return minimal content to allow processing to continue
            return f"PDF extraction error: {str(e)}", {
                "pages": 0,
                "text_by_page": {},
                "error": str(e),
                "format": "pdf"
            }
    
    @staticmethod
    def _extract_docx(file_content: bytes) -> Tuple[str, dict]:
        """Extract text from DOCX with structure"""
        try:
            doc = DocxDocument(io.BytesIO(file_content))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            full_text = "\n".join(text_parts)
            
            return full_text, {
                "paragraphs": len(doc.paragraphs),
                "sections": len(doc.sections),
                "tables": len(doc.tables),
                "format": "docx"
            }
        except Exception as e:
            logger.error(f"DOCX extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_pptx(file_content: bytes) -> Tuple[str, dict]:
        """Extract text from PPTX with slide markers"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"--- SLIDE {slide_num} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)
                
                if len(slide_parts) > 1:  # If there's content
                    slides_text.append("\n".join(slide_parts))
            
            full_text = "\n\n".join(slides_text)
            
            return full_text, {
                "slides": len(prs.slides),
                "format": "pptx"
            }
        except Exception as e:
            logger.error(f"PPTX extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_txt(file_content: bytes) -> Tuple[str, dict]:
        """Extract text from plain text file"""
        try:
            text = file_content.decode("utf-8", errors="ignore")
            lines = text.split("\n")
            
            return text, {
                "lines": len(lines),
                "format": "txt"
            }
        except Exception as e:
            logger.error(f"TXT extraction error: {e}")
            raise
    
    @staticmethod
    def _extract_csv(file_content: bytes) -> Tuple[str, dict]:
        """Extract and format CSV content"""
        try:
            df = pd.read_csv(io.BytesIO(file_content))
            
            # Convert to formatted text
            text = df.to_string()
            
            return text, {
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": list(df.columns),
                "format": "csv"
            }
        except Exception as e:
            logger.error(f"CSV extraction error: {e}")
            raise


class TextChunker:
    """Intelligently chunk text for embeddings and retrieval"""
    
    @staticmethod
    def chunk_text(
        text: str,
        chunk_size: int = 512,
        overlap: int = 50,
        preserve_pages: bool = True
    ) -> List[dict]:
        """Split text into overlapping chunks with metadata
        
        Args:
            text: Full document text
            chunk_size: Target chunk size in characters
            overlap: Overlap between chunks in characters
            preserve_pages: Keep page markers in chunks
        
        Returns:
            List of {content, page, index}
        """
        chunks = []
        
        # Split by sentences for better context
        sentences = TextChunker._smart_split(text)
        
        current_chunk = ""
        current_page = 1
        chunk_index = 0
        
        for sentence in sentences:
            # Track page numbers
            if "--- PAGE" in sentence or "--- SLIDE" in sentence:
                try:
                    page_str = re.search(r"--- (?:PAGE|SLIDE) (\d+)", sentence).group(1)
                    current_page = int(page_str)
                except:
                    pass
            
            # Check if adding this sentence would exceed chunk_size
            test_chunk = current_chunk + " " + sentence if current_chunk else sentence
            
            if len(test_chunk) <= chunk_size:
                current_chunk = test_chunk.strip()
            else:
                # Save current chunk if it has content
                if current_chunk.strip():
                    chunks.append({
                        "content": current_chunk.strip(),
                        "page": current_page,
                        "index": chunk_index,
                        "size": len(current_chunk)
                    })
                    chunk_index += 1
                
                # Start new chunk with overlap
                if overlap > 0 and len(current_chunk) > overlap:
                    # Keep last `overlap` characters as context
                    current_chunk = current_chunk[-overlap:] + " " + sentence
                else:
                    current_chunk = sentence
        
        # Add final chunk
        if current_chunk.strip():
            chunks.append({
                "content": current_chunk.strip(),
                "page": current_page,
                "index": chunk_index,
                "size": len(current_chunk)
            })
        
        logger.info(f"Created {len(chunks)} chunks with average size {sum(c['size'] for c in chunks) / len(chunks):.0f}")
        return chunks
    
    @staticmethod
    def _smart_split(text: str) -> List[str]:
        """Split text intelligently by sentences while preserving structure"""
        # First split by double newlines (paragraphs)
        paragraphs = text.split("\n\n")
        sentences = []
        
        for para in paragraphs:
            # Split paragraph into sentences
            # Handle common abbreviations
            para = re.sub(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', '|||SENT_END|||', para)
            
            para_sentences = para.split('|||SENT_END|||')
            
            for sent in para_sentences:
                sent = sent.strip()
                if sent:
                    sentences.append(sent)
        
        return sentences


class EmbeddingGenerator:
    """Generate embeddings using BGE-M3 or other models"""
    
    def __init__(self, model_name: str = "BAAI/bge-m3"):
        self.model_name = model_name
        self.model = None
        self.dimension = 1024
        self._model_loaded = False
        logger.info(f"Initialized embedding generator for {model_name} (lazy loading)")
    
    def _load_model(self):
        """Lazy load the embedding model"""
        if self._model_loaded:
            return
        
        try:
            import os
            # Disable symlinks for HuggingFace cache on Windows
            os.environ['HF_HUB_DISABLE_SYMLINKS_WARNING'] = '1'
            
            from sentence_transformers import SentenceTransformer
            logger.info(f"Loading embedding model: {self.model_name}")
            
            # Use trust_remote_code=True to handle model configuration
            self.model = SentenceTransformer(
                self.model_name,
                trust_remote_code=True,
                device='cpu'  # Use CPU to avoid CUDA issues
            )
            
            # Get actual dimension
            dummy_embedding = self.model.encode(["test"])
            self.dimension = dummy_embedding.shape[1]
            self._model_loaded = True
            logger.info(f"Loaded {self.model_name} with dimension {self.dimension}")
        
        except ImportError:
            logger.warning("sentence-transformers not available, using dummy embeddings")
            self._model_loaded = True
        except Exception as e:
            logger.warning(f"Failed to load embedding model: {e}", exc_info=True)
            # Set model_loaded to True even on error to prevent retry loops
            self._model_loaded = True
    
    def generate(self, texts: List[str], batch_size: int = 32) -> List[List[float]]:
        """Generate embeddings for multiple texts
        
        Args:
            texts: List of texts to embed
            batch_size: Process texts in batches for efficiency
        
        Returns:
            List of embedding vectors
        """
        self._load_model()  # Ensure model is loaded
        
        if not self.model:
            logger.warning("Using dummy embeddings (model not loaded)")
            return [[0.0] * self.dimension for _ in texts]
        
        try:
            embeddings = self.model.encode(
                texts,
                batch_size=batch_size,
                show_progress_bar=False,
                normalize_embeddings=True
            )
            return embeddings.tolist()
        
        except Exception as e:
            logger.error(f"Embedding generation error: {e}")
            return [[0.0] * self.dimension for _ in texts]
    
    def generate_single(self, text: str) -> List[float]:
        """Generate embedding for single text"""
        return self.generate([text], batch_size=1)[0]
    
    def get_dimension(self) -> int:
        """Get embedding dimension"""
        return self.dimension
